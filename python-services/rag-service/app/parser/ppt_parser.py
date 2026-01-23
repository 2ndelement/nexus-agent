"""
app/parser/ppt_parser.py — PPT 解析器
"""
import logging
from pathlib import Path

from app.parser.base import BaseParser, ParseResult

logger = logging.getLogger(__name__)


class PPTParser(BaseParser):
    """PPT 文档解析器"""
    
    def supports(self, file_path: str) -> bool:
        """支持 PPT 文件"""
        return Path(file_path).suffix.lower() in ['.pptx', '.ppt']
    
    def parse(self, file_path: str) -> ParseResult:
        """解析 PPT 文档"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_parts = []
            
            # 提取标题
            title = prs.core_properties.title or Path(file_path).stem
            
            metadata = {
                "slide_count": len(prs.slides),
                "source_file": file_path,
                "file_type": "pptx"
            }
            
            # 遍历所有幻灯片
            for slide_num, slide in enumerate(prs.slides):
                slide_text = [f"[幻灯片 {slide_num + 1}]"]
                
                # 提取标题
                if slide.shapes.title:
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        slide_text.append(f"## {title_text}")
                
                # 提取所有文本框
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # 排除标题（已单独处理）
                        if shape != slide.shapes.title:
                            slide_text.append(shape.text.strip())
                
                # 提取表格
                for table in shape.table if hasattr(shape, 'table') else []:
                    slide_text.append("\n表格:")
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        slide_text.append(row_text)
                
                if len(slide_text) > 1:  # 有内容
                    text_parts.append("\n".join(slide_text))
            
            content = "\n\n".join(text_parts)
            
            return ParseResult(
                content=content,
                title=title,
                metadata=metadata,
                success=True
            )
            
        except ImportError:
            return ParseResult(
                content="",
                title=Path(file_path).stem,
                metadata={"file_type": "pptx"},
                success=False,
                error_message="python-pptx not installed: pip install python-pptx"
            )
        except Exception as e:
            logger.error(f"PPT parsing error: {e}")
            return ParseResult(
                content="",
                title=Path(file_path).stem,
                metadata={"file_type": "pptx"},
                success=False,
                error_message=str(e)
            )
