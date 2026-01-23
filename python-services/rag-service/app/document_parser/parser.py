"""
app/document_parser/parser.py — 多格式文档解析器

支持的格式：
- PDF (.pdf)
- Word (.docx)
- PowerPoint (.pptx)
- TXT (.txt)
- Markdown (.md)
- HTML (.html)
"""
from __future__ import annotations

import logging
import re
from abc import ABC, abstractmethod
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


@dataclass
class ParsedDocument:
    """解析后的文档"""
    title: str                    # 文档标题
    content: str                  # 纯文本内容
    metadata: dict               # 元数据（页码、格式等）
    sections: list[Section]       # 结构化段落
    

@dataclass  
class Section:
    """文档章节"""
    level: int           # 层级 (1=标题1, 2=标题2, 0=正文)
    title: Optional[str] # 标题（正文为None）
    content: str         # 内容
    start_line: int      # 起始行号
    end_line: int        # 结束行号


class BaseParser(ABC):
    """解析器基类"""
    
    @abstractmethod
    def parse(self, file_path: str) -> ParsedDocument:
        """解析文档"""
        pass
    
    @property
    @abstractmethod
    def supported_extensions(self) -> list[str]:
        """支持的扩展名"""
        pass


class TextParser(BaseParser):
    """纯文本解析器"""
    
    @property
    def supported_extensions(self) -> list[str]:
        return [".txt", ".md", ".markdown"]
    
    def parse(self, file_path: str) -> ParsedDocument:
        path = Path(file_path)
        
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
        
        # 提取标题（从文件名或第一行）
        title = path.stem
        if content.strip():
            first_line = content.split("\n")[0].strip()
            if first_line.startswith("#"):
                title = first_line.lstrip("#").strip()
        
        # 解析章节结构
        sections = self._parse_sections(content)
        
        return ParsedDocument(
            title=title,
            content=content,
            metadata={
                "format": path.suffix,
                "file_name": path.name,
                "char_count": len(content),
            },
            sections=sections,
        )
    
    def _parse_sections(self, content: str) -> list[Section]:
        """解析章节结构"""
        sections = []
        lines = content.split("\n")
        
        i = 0
        current_level = 0
        current_title = None
        section_start = 0
        section_content_lines = []
        
        while i < len(lines):
            line = lines[i]
            
            # 检测标题 (# ## ###)
            match = re.match(r"^(#{1,6})\s+(.+)$", line)
            if match:
                # 保存上一个段落
                if section_content_lines:
                    sections.append(Section(
                        level=current_level,
                        title=current_title,
                        content="\n".join(section_content_lines).strip(),
                        start_line=section_start,
                        end_line=i - 1,
                    ))
                    section_content_lines = []
                
                # 新标题
                current_level = len(match.group(1))
                current_title = match.group(2).strip()
                section_start = i
            
            else:
                # 正文内容
                if current_level == 0:
                    current_level = 0
                section_content_lines.append(line)
            
            i += 1
        
        # 保存最后一个段落
        if section_content_lines:
            sections.append(Section(
                level=current_level,
                title=current_title,
                content="\n".join(section_content_lines).strip(),
                start_line=section_start,
                end_line=len(lines) - 1,
            ))
        
        return sections


class PDFParser(BaseParser):
    """PDF 解析器"""
    
    @property
    def supported_extensions(self) -> list[str]:
        return [".pdf"]
    
    def parse(self, file_path: str) -> ParsedDocument:
        try:
            import pymupdf  # PyMuPDF
        except ImportError:
            logger.warning("PyMuPDF 未安装，使用 fallback")
            return self._fallback_parse(file_path)
        
        path = Path(file_path)
        doc = pymupdf.open(file_path)
        
        title = path.stem
        all_text = ""
        sections = []
        
        for page_num, page in enumerate(doc):
            text = page.get_text()
            all_text += text + "\n"
            
            # 简单按页分段
            if text.strip():
                sections.append(Section(
                    level=0,
                    title=f"第{page_num + 1}页",
                    content=text.strip(),
                    start_line=page_num,
                    end_line=page_num,
                ))
        
        doc.close()
        
        return ParsedDocument(
            title=title,
            content=all_text.strip(),
            metadata={
                "format": ".pdf",
                "file_name": path.name,
                "page_count": len(doc),
                "char_count": len(all_text),
            },
            sections=sections,
        )
    
    def _fallback_parse(self, file_path: str) -> ParsedDocument:
        """Fallback: 使用 pdfplumber"""
        try:
            import pdfplumber
            
            path = Path(file_path)
            all_text = ""
            
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        all_text += text + "\n"
            
            return ParsedDocument(
                title=path.stem,
                content=all_text.strip(),
                metadata={"format": ".pdf", "file_name": path.name},
                sections=[],
            )
        except Exception as e:
            logger.error("PDF 解析失败: %s", e)
            return ParsedDocument(
                title=Path(file_path).stem,
                content="",
                metadata={"error": str(e)},
                sections=[],
            )


class WordParser(BaseParser):
    """Word 文档解析器"""
    
    @property
    def supported_extensions(self) -> list[str]:
        return [".docx"]
    
    def parse(self, file_path: str) -> ParsedDocument:
        try:
            from docx import Document
        except ImportError:
            logger.warning("python-docx 未安装")
            return self._empty_doc(file_path)
        
        path = Path(file_path)
        doc = Document(file_path)
        
        title = path.stem
        all_text = ""
        sections = []
        
        # 提取标题
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            # 检测标题样式
            style_name = para.style.name if para.style else ""
            if "Heading" in style_name or para.runs and any(run.bold for run in para.runs if run.text):
                level = int(style_name.replace("Heading ", "")) if "Heading" in style_name else 1
                sections.append(Section(
                    level=level,
                    title=text,
                    content="",
                    start_line=len(all_text.split("\n")),
                    end_line=0,
                ))
            else:
                all_text += text + "\n"
        
        return ParsedDocument(
            title=title,
            content=all_text.strip(),
            metadata={
                "format": ".docx",
                "file_name": path.name,
                "paragraph_count": len(doc.paragraphs),
            },
            sections=sections,
        )
    
    def _empty_doc(self, file_path: str) -> ParsedDocument:
        return ParsedDocument(
            title=Path(file_path).stem,
            content="",
            metadata={"error": "python-docx 未安装"},
            sections=[],
        )


class PPTParser(BaseParser):
    """PowerPoint 解析器"""
    
    @property
    def supported_extensions(self) -> list[str]:
        return [".pptx"]
    
    def parse(self, file_path: str) -> ParsedDocument:
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx 未安装")
            return self._empty_doc(file_path)
        
        path = Path(file_path)
        prs = Presentation(file_path)
        
        title = path.stem
        all_text = ""
        sections = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = ""
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text += shape.text + "\n"
            
            if slide_text.strip():
                all_text += f"=== 第{slide_num + 1}页 ===\n"
                all_text += slide_text.strip() + "\n\n"
                
                sections.append(Section(
                    level=0,
                    title=f"第{slide_num + 1}页",
                    content=slide_text.strip(),
                    start_line=slide_num,
                    end_line=slide_num,
                ))
        
        return ParsedDocument(
            title=title,
            content=all_text.strip(),
            metadata={
                "format": ".pptx",
                "file_name": path.name,
                "slide_count": len(prs.slides),
            },
            sections=sections,
        )
    
    def _empty_doc(self, file_path: str) -> ParsedDocument:
        return ParsedDocument(
            title=Path(file_path).stem,
            content="",
            metadata={"error": "python-pptx 未安装"},
            sections=[],
        )


class DocumentParser:
    """文档解析器统一入口"""
    
    def __init__(self):
        self.parsers: dict[str, BaseParser] = {}
        
        # 注册所有解析器
        for parser in [TextParser(), PDFParser(), WordParser(), PPTParser()]:
            for ext in parser.supported_extensions:
                self.parsers[ext.lower()] = parser
        
        logger.info("文档解析器已加载，支持格式: %s", list(self.parsers.keys()))
    
    def parse(self, file_path: str) -> ParsedDocument:
        """解析文档"""
        ext = Path(file_path).suffix.lower()
        
        parser = self.parsers.get(ext)
        if not parser:
            raise ValueError(f"不支持的文档格式: {ext}")
        
        return parser.parse(file_path)
    
    def parse_content(self, content: str, format: str = "txt") -> ParsedDocument:
        """从内容解析（用于API上传）"""
        # 简单处理：直接返回文本
        return ParsedDocument(
            title="文档",
            content=content,
            metadata={"format": format},
            sections=[Section(level=0, title=None, content=content, start_line=0, end_line=0)],
        )
